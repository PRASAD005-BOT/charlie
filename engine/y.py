import google.generativeai as genai  # Ensure you have the Google Gemini API client library installed
from flask import Flask, render_template, request, send_file
from flask_limiter import Limiter
import os
import random
import re
from pptx import Presentation

# Initialize the Flask application
app = Flask(__name__)

# Rate limiting
limiter = Limiter(app, default_limits=["10 per day"])

# Configure Google Generative AI API
api_key = os.getenv('---------------------------')  # Replace with your environment variable for the API key
genai.configure(api_key=api_key)
model = genai.GenerativeModel("gemini-1.5-flash")

def create_presentation(text_file, design_number, ppt_name):
    """Create a PowerPoint presentation based on generated text."""
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0
    last_slide_layout_index = -1
    content = ""

    with open(text_file, 'r', encoding='utf-8') as f:
        for line in f:
            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = header
            elif line.startswith('#Slide:'):
                if slide_count > 0:  # Add previous content to the slide
                    slide = prs.slides.add_slide(prs.slide_layouts[last_slide_layout_index])
                    body_shape = slide.shapes.placeholders[1]
                    body_shape.text_frame.text = content
                content = ""
                slide_count += 1
                last_slide_layout_index = random.choice([1, 7, 8])  # Random layout
            elif line.startswith('#Header:'):
                continue  # Headers are not directly used
            elif line.startswith('#Content:'):
                content = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()

    # Save the presentation
    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')

    # Clean up the temporary text file
    if os.path.exists(text_file):
        os.remove(text_file)

    return f"{request.host_url}GeneratedPresentations/{ppt_name}.pptx"

@app.route('/GeneratedPresentations/<path:path>')
def send_generated_presentation(path):
    """Serve the generated presentation file."""
    return send_file(f'GeneratedPresentations/{path}', as_attachment=True)

@app.route("/powerpoint")
def powerpoint():
    """Render the PowerPoint input page."""
    return render_template("powerpoint.html", charset="utf-8")

@app.route("/")
def home():
    """Render the home page."""
    return render_template("powerpoint.html", charset="utf-8")

@app.route("/get")
@limiter.limit("10 per day")
def get_bot_response():
    """Handle the user's request and generate a PowerPoint presentation."""
    user_text = request.args.get("msg")
    user_text = re.sub(r'[^\w\s.\-\(\)]', '', user_text).strip()
    design_number = int(user_text[-1]) if user_text[-1].isdigit() else 1
    user_input = user_text[:-2] if user_text[-1].isdigit() else user_text

    # Ensure the design number is valid
    design_number = min(max(design_number, 1), 7)  # Clamp to [1, 7]

    # Generate a filename based on the user's input
    filename = re.sub(r'\s+', '_', user_input)  # Replace spaces with underscores

    try:
        # Generate presentation text using Google Gemini AI's API
        presentation_text = model.generate(
            prompt=f"The user wants a presentation about {user_input}.",
            temperature=0.5,
            max_tokens=500  # Adjust based on your needs
        )['output']  # Update according to the response structure

        with open(f'Cache/{filename}.txt', 'w', encoding='utf-8') as f:
            f.write(presentation_text)

        # Create the PowerPoint presentation
        ppt_link = create_presentation(f'Cache/{filename}.txt', design_number, filename)
        return ppt_link
    except Exception as e:
        return f"An error occurred: {str(e)}", 500

if __name__ == '__main__':
    # Set debug=True for development
    app.run(debug=True)
