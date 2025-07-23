import google.generativeai as genai  # Import the Generative AI library
import os
import eel
from pptx import Presentation
from pptx.util import Pt
import speech_recognition as sr  # For voice recognition
import pyttsx3 
from pptx.util import Cm, Pt # For text-to-speech
eel.init('web')

def DisplayMessage(message: str):
    """Display the message in a pop-up window using eel."""
    eel.DisplayMessage(message)
model = genai.GenerativeModel("gemini-1.5-flash")

# Initialize Text-to-Speech engine
tts_engine = pyttsx3.init()

# Create directories if they don't exist
os.makedirs('GeneratedPresentations', exist_ok=True)
def receiveLocalStorageData(data):
    global stored_chatbot_keys
    stored_chatbot_keys = {}
    chatbot_keys = data.get('ChatbotKeys', [])
    for key in chatbot_keys:
        if key.get('keyName') == 'chatbotkey':
            api_key = key.get('apiKey')
            if api_key:
                stored_chatbot_keys['chatbotkey'] = api_key

                # Configure genai dynamically
                print(f"Configuring genai with API key: ")
                try:
                    genai.configure(api_key=api_key)
                except Exception as e:
                    print(f"Error configuring genai: {e}")
def listen_for_commands(prompt):
    """Capture voice input and convert it to text using speech recognition."""
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        print(prompt)
        tts_engine.say(prompt)
        tts_engine.runAndWait()
        recognizer.adjust_for_ambient_noise(source)
        audio = recognizer.listen(source)

        try:
            command = recognizer.recognize_google(audio)
            eel.DisplayMessage(f"Recognized command: {command}")
            print(f"Recognized command: {command}")
            return command
        except sr.UnknownValueError:
            eel.DisplayMessage("Sorry, I could not understand the audio.")
            print("Sorry, I could not understand the audio.")
            return None
        except sr.RequestError as e:
            eel.DisplayMessage(f"Could not request results from Google Speech Recognition service; {e}")
            print(f"Could not request results from Google Speech Recognition service; {e}")
            return None

def chat_with_gemini(extracted_text, query):
    """Interact with Gemini AI to generate text based on extracted text and a query."""
    extracted_text = extracted_text.strip() if extracted_text else "No text provided."
    query = query.strip() if query else "No query provided."

    response = model.generate_content([extracted_text, query])

    if hasattr(response, 'text'):
        generated_text = response.text.strip().replace('*', '')
        return generated_text
    else:
        return "No response received."
def generate_slide_content(topic, slide_index):
    """
    Generate unique content for a slide about the specified topic.
    Each slide will have a slightly different focus based on the index.
    """
    # Modify extracted_text and query based on the slide index
    extracted_text = f"Discuss key aspects of {topic} focusing on detail {slide_index + 1}."
    query = f"Provide detailed information about {topic}, focusing on aspect {slide_index + 1} of its importance."

    # Generate content using Gemini AI
    content = chat_with_gemini(extracted_text, query)
    return content
def create_presentation(topic, slide_contents, design_path=None):
    """Create a PowerPoint presentation with slides fitting specified dimensions."""
    prs = Presentation(design_path) if design_path else Presentation()  # Use provided design or blank presentation

    # Define the blank slide layout
    blank_layout = prs.slide_layouts[6]  # Typically, index 6 is the blank layout

    for i, content in enumerate(slide_contents):
        slide = prs.slides.add_slide(blank_layout)  # Add a blank slide

        # Add a title shape
        title_shape = slide.shapes.add_textbox(left=Cm(1), top=Cm(1), width=Cm(30), height=Cm(2))  # Title box
        title_frame = title_shape.text_frame
        title_frame.text = f"{topic} - Slide {i + 1}"

        # Add content shape with specified box dimensions
        content_shape = slide.shapes.add_textbox(left=Cm(1), top=Cm(4), width=Cm(31.57), height=Cm(16.86))
        content_frame = content_shape.text_frame

        # Limit content length to fit within the box
        max_length = 2000  # Define a reasonable limit
        if len(content) > max_length:
            content = content[:max_length] + '...'  # Truncate text with ellipsis

        # Set content text
        content_frame.text = content

        # Adjust font size and style to fit within the box
        for paragraph in content_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)  # Adjust font size as needed
                run.font.name = "Arial"  # Set preferred font

    # Save the presentation
    ppt_name = f"{topic.replace(' ', '_')}_presentation.pptx"
    prs.save(f'GeneratedPresentations/{ppt_name}')

    return f'GeneratedPresentations/{ppt_name}'

def listen_for_valid_slide_number():
    """Listen and continuously prompt the user for a valid number of slides (1-100)."""
    while True:
        eel.DisplayMessage("How many slides would you like? Please say a number between 1 and 100.")
        num_slides = listen_for_commands("How many slides would you like? Please say a number between 1 and 100.")

        if num_slides is not None:
            # Attempt to extract numbers even if the user says "1 slide" or "16 slides"
            num_slides = ''.join(filter(str.isdigit, num_slides))  # Extract only digits from the spoken input
            
            if num_slides.isdigit():  # Check if the remaining part is a valid digit sequence
                num_slides_int = int(num_slides)  # Convert spoken text to an integer
                
                if 1 <= num_slides_int <= 100:
                    return num_slides_int
                else:
                    eel.DisplayMessage("How many slides would you like? Please say a number between 1 and 100.")
                    tts_engine.say("Please provide a number between 1 and 100.")
            else:
                eel.DisplayMessage("Invalid input. Please say a valid number between 1 and 100.")
                tts_engine.say("Invalid input. Please say a valid number between 1 and 100.")
        else:
            eel.DisplayMessage("Sorry, I didn't understand that. Could you repeat?")
            tts_engine.say("Sorry, I didn't understand that. Could you repeat?")
        
        tts_engine.runAndWait()


def listen_for_design_file():
    """Prompt the user to select a design file for the PowerPoint."""
    available_designs = os.listdir(r'D:\pytho\jarvis\powerpoint\Designs')  # List available designs
    if not available_designs:
        eel.DisplayMessage("No design files available. Proceeding with a blank presentation.")
        tts_engine.say("No design files available. Proceeding with a blank presentation.")
        tts_engine.runAndWait()
        return None
    eel.DisplayMessage("Please select a design for the PowerPoint. Available designs are:")
    tts_engine.say("Please select a design for the PowerPoint. Available designs are:")
    for design in available_designs:
        tts_engine.say(design)
    tts_engine.runAndWait()

    while True:
        eel.DisplayMessage("Which design would you like? Please say the design number.")
        selected_design = listen_for_commands("Which design would you like? Please say the design number.")
        if selected_design is not None:
            # Map user input to the design file
            design_map = {
                "design 01": "design-1.pptx",
                "design 02": "design-2.pptx",
                "design 03": "design-3.pptx",
                "design 04": "design-4.pptx",
                "design 05": "design-5.pptx",
                "design 06": "design-6.pptx",
                "design 07": "design-7.pptx",
                "design blue": "Design-Blue.pptx"

            }

            if selected_design.lower() in design_map:
                return design_map[selected_design.lower()]
            else:
                eel.DisplayMessage("Invalid design. Please select from the available designs.")
                tts_engine.say("Invalid design. Please select from the available designs.")
                tts_engine.runAndWait()
def powerpoint():
    """Main function to capture voice input, generate content, and create a presentation."""
    # Listen for the topic
    eel.DisplayMessage("Please provide the topic for your presentation.")
    topic = listen_for_commands("Please provide the topic for your presentation.")
    if topic is None:
        eel.DisplayMessage("Sorry, I could not understand the topic.")
        tts_engine.say("Sorry, I could not understand the topic.")
        tts_engine.runAndWait()
        return

    # Get a valid number of slides
    num_slides = listen_for_valid_slide_number()

    # Listen for the design file
    design_file = listen_for_design_file()

    if design_file:
        design_path = os.path.join(r'D:\pytho\jarvis\powerpoint\Designs', design_file)  # Use the selected design
    else:
        design_path = None  # Use default blank presentation if no design is selected

    # Inform user that PowerPoint generation is starting
    eel.DisplayMessage("Your PowerPoint is being generated. Please wait a moment.")
    tts_engine.say("Your PowerPoint is being generated. Please wait a moment.")
    tts_engine.runAndWait()

    # Generate unique content for each slide
    slide_contents = []
    for slide_index in range(num_slides):
        content = generate_slide_content(topic, slide_index)
        slide_contents.append(content)

    # Create the PowerPoint presentation
    ppt_link = create_presentation(topic, slide_contents, design_path)

    # Provide feedback to the user
    eel.DisplayMessage(f"Your presentation has been created. It is saved as {ppt_link}. on your desktop, Boss!")
    tts_engine.say(f"Your presentation has been created. It is saved on your desktop, Boss!")
    tts_engine.runAndWait()


